"""rag_document_loader.py

Minimal, high-performance, modular document loader for RAG pipelines.
Supports: .txt, .md, .csv, .html, .pdf, .docx, .pptx, images (OCR), directories, .zip

Design:
- Pluggable handlers via a registry
- Async-friendly: uses ThreadPoolExecutor for blocking I/O
- Graceful degradation when optional dependencies are missing
- Minimal external deps; optional: PyPDF2, pdfminer.six, python-docx, python-pptx, pillow, pytesseract, beautifulsoup4

Usage example at bottom.
"""

from __future__ import annotations
import asyncio
import concurrent.futures
import csv
import io
import os
import re
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Dict, Iterable, List, Optional, Tuple

# Optional imports (handled gracefully)
try:
    from PyPDF2 import PdfReader  # pip install PyPDF2
except Exception:
    PdfReader = None

try:
    import docx  # pip install python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation  # pip install python-pptx
except Exception:
    Presentation = None

try:
    from PIL import Image  # pip install pillow
except Exception:
    Image = None

try:
    import pytesseract  # pip install pytesseract (requires tesseract installed on system)
except Exception:
    pytesseract = None

try:
    from bs4 import BeautifulSoup  # pip install beautifulsoup4
except Exception:
    BeautifulSoup = None


# # Simple text splitter (can be replaced/injected by caller)
def default_text_splitter(text: str, chunk_size: int = 1024, overlap: int = 128) -> List[str]:
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start = max(end - overlap, end)
        if start == end and end >= len(text):
            break
    return chunks


@dataclass
class Document:
    source: str
    content: str
    metadata: Dict[str, str]


class DocumentLoader:
    """Registry-based, async-friendly document loader."""

    def __init__(self, max_workers: int = None):
        self.handlers: Dict[str, Callable[[Path], str]] = {}
        self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
        self._register_default_handlers()

    # Public API -------------------------------------------------
    async def load(self, paths: Iterable[str], concurrency: int = 8) -> List[Document]:
        """Load multiple paths (files or directories). Returns list of Document objects."""
        loop = asyncio.get_running_loop()
        tasks = []
        sem = asyncio.Semaphore(concurrency)

        async def _load_path(p: str):
            async with sem:
                return await loop.run_in_executor(self.executor, self._load_sync, p)

        for p in paths:
            tasks.append(asyncio.create_task(_load_path(p)))

        results = await asyncio.gather(*tasks)
        # flatten
        docs: List[Document] = []
        for r in results:
            if r:
                docs.extend(r)
        return docs

    def register_handler(self, suffix: str, func: Callable[[Path], Optional[List[Document]]]):
        """Register custom handler for file suffix (e.g. '.xyz'). Handler returns list of Documents or None."""
        self.handlers[suffix.lower()] = func

    # Internal helpers -------------------------------------------
    def _load_sync(self, p: str) -> List[Document]:
        pth = Path(p)
        if not pth.exists():
            return []
        if pth.is_dir():
            return self._load_directory(pth)
        if pth.suffix.lower() == ".zip":
            return self._load_zip(pth)
        handler = self.handlers.get(pth.suffix.lower())
        if handler:
            out = handler(pth)
            return out or []
        # fallback: try text
        return [Document(source=str(pth), content=pth.read_text(encoding="utf-8", errors="ignore"), metadata={"type": "text"})]

    def _load_directory(self, dirpath: Path) -> List[Document]:
        docs: List[Document] = []
        for root, _dirs, files in os.walk(dirpath, followlinks=True):
            for fname in files:
                full = Path(root) / fname
                docs.extend(self._load_sync(str(full)))
        return docs

    def _load_zip(self, zippath: Path) -> List[Document]:
        docs: List[Document] = []
        with zipfile.ZipFile(zippath, 'r') as zf:
            with tempfile.TemporaryDirectory() as td:
                zf.extractall(td)
                docs.extend(self._load_directory(Path(td)))
        return docs

    # Handlers --------------------------------------------------
    def _register_default_handlers(self):
        self.register_handler('.txt', self._handle_text_file)
        self.register_handler('.md', self._handle_text_file)
        self.register_handler('.csv', self._handle_csv)
        self.register_handler('.html', self._handle_html)
        self.register_handler('.htm', self._handle_html)
        self.register_handler('.pdf', self._handle_pdf)
        self.register_handler('.docx', self._handle_docx)
        self.register_handler('.pptx', self._handle_pptx)
        self.register_handler('.png', self._handle_image)
        self.register_handler('.jpg', self._handle_image)
        self.register_handler('.jpeg', self._handle_image)
        self.register_handler('.gif', self._handle_image)

    def _handle_text_file(self, path: Path) -> List[Document]:
        txt = path.read_text(encoding='utf-8', errors='ignore')
        return [Document(source=str(path), content=txt, metadata={"type": "text"})]

    def _handle_csv(self, path: Path) -> List[Document]:
        docs: List[Document] = []
        with path.open('r', encoding='utf-8', errors='ignore', newline='') as fh:
            reader = csv.reader(fh)
            rows = []
            for row in reader:
                rows.append(', '.join(row))
            text = '\n'.join(rows)
            docs.append(Document(source=str(path), content=text, metadata={"type": "csv", "rows": str(len(rows))}))
        return docs

    def _handle_html(self, path: Path) -> List[Document]:
        raw = path.read_text(encoding='utf-8', errors='ignore')
        if BeautifulSoup:
            soup = BeautifulSoup(raw, 'html.parser')
            text = soup.get_text(separator='\n', strip=True)
        else:
            # fallback naive strip
            text = re.sub(r'<[^>]+>', '', raw)
        return [Document(source=str(path), content=text, metadata={"type": "html"})]

    def _handle_pdf(self, path: Path) -> List[Document]:
        if PdfReader:
            try:
                reader = PdfReader(str(path))
                texts = []
                for p in reader.pages:
                    try:
                        texts.append(p.extract_text() or '')
                    except Exception:
                        # per-page fallback
                        texts.append('')
                combined = "\n".join(texts)
                return [Document(source=str(path), content=combined, metadata={"type": "pdf", "pages": str(len(reader.pages))})]
            except Exception:
                pass
        # Generic binary fallback: try reading as text (often garbage)
        try:
            raw = path.read_text(encoding='utf-8', errors='ignore')
            return [Document(source=str(path), content=raw, metadata={"type": "pdf.fallback"})]
        except Exception:
            return []

    def _handle_docx(self, path: Path) -> List[Document]:
        if docx is None:
            return [Document(source=str(path), content='', metadata={"type": "docx", "warning": "python-docx not installed"})]
        doc = docx.Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text]
        return [Document(source=str(path), content='\n'.join(paras), metadata={"type": "docx", "paragraphs": str(len(paras))})]

    def _handle_pptx(self, path: Path) -> List[Document]:
        if Presentation is None:
            return [Document(source=str(path), content='', metadata={"type": "pptx", "warning": "python-pptx not installed"})]
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    t = shape.text
                    if t:
                        texts.append(t)
        return [Document(source=str(path), content='\n'.join(texts), metadata={"type": "pptx", "slides": str(len(prs.slides))})]

    def _handle_image(self, path: Path) -> List[Document]:
        if Image is None:
            return [Document(source=str(path), content='', metadata={"type": "image", "warning": "pillow not installed"})]
        if pytesseract is None:
            return [Document(source=str(path), content='', metadata={"type": "image", "warning": "pytesseract not installed"})]
        try:
            img = Image.open(path)
            text = pytesseract.image_to_string(img)
            return [Document(source=str(path), content=text, metadata={"type": "image.ocr"})]
        except Exception:
            return [Document(source=str(path), content='', metadata={"type": "image", "warning": "ocr_failed"})]


# Small utility to convert Documents into text chunks for RAG
async def document_load(loader: DocumentLoader, paths: Iterable[str], chunk_size: int = 1024, overlap: int = 128) -> List[Tuple[str, str, Dict[str, str]]]:
    docs = await loader.load(paths)
    out: List[Tuple[str, str, Dict[str, str]]] = []
    for d in docs:
        chunks = default_text_splitter(d.content, chunk_size=chunk_size, overlap=overlap)
        for i, c in enumerate(chunks):
            meta = dict(d.metadata)
            meta.update({"chunk_index": str(i)})
            out.append((d.source, c, meta))
    return out 


# Example usage / minimal test --------------------------------
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Load documents for RAG pipelines")
    parser.add_argument('paths', nargs='*', default=['artifacts/'], help='Files or directories to load (default: artifacts/)')
    parser.add_argument('--chunk', type=int, default=1024)
    parser.add_argument('--overlap', type=int, default=128)
    args = parser.parse_args()

    async def main():
        loader = DocumentLoader()
        # Example: register custom handler (demonstration)
        # loader.register_handler('.foo', lambda p: [Document(source=str(p), content='custom', metadata={})])
        items = await document_load(loader, args.paths, chunk_size=args.chunk, overlap=args.overlap)
        # chunks = default_text_splitter(d.content, chunk_size=args.chunk, overlap=args.overlap)
        print(f'Loaded {len(items)} chunks from {len(args.paths)} input paths')
        # print first 3 chunks
        for src, txt, meta in items[:3]:
            print('---')
            print('source:', src)
            print('meta:', meta)
            print('text:', (txt[:300] + '...') if len(txt) > 300 else txt)

    asyncio.run(main())
